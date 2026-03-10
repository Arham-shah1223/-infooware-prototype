import os
from moviepy import ImageClip, concatenate_videoclips
from pptx import Presentation
from PIL import Image, ImageDraw


def slides_to_images(pptx_path, temp_folder):

    prs = Presentation(pptx_path)

    images = []

    if not os.path.exists(temp_folder):
        os.makedirs(temp_folder)

    for i, slide in enumerate(prs.slides):

        img_path = os.path.join(temp_folder, f"slide_{i}.png")

        img = Image.new("RGB", (1280, 720), "white")
        draw = ImageDraw.Draw(img)

        title = slide.shapes.title.text if slide.shapes.title else ""

        # draw title
        draw.text((100, 80), title, fill="black")

        content_lines = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text.split("\n"):
                    if paragraph.strip() != "":
                        content_lines.append(paragraph)

        y = 200

        for line in content_lines:
            draw.text((100, y), line, fill="black")
            y += 40

        img.save(img_path)

        images.append(img_path)

    return images


def create_video(pptx_path, output_video):

    temp_folder = "temp_slides"

    slide_images = slides_to_images(pptx_path, temp_folder)

    clips = []

    for img in slide_images:

        clip = ImageClip(img).with_duration(5)

        clips.append(clip)

    video = concatenate_videoclips(clips)

    video.write_videofile(output_video, fps=24)