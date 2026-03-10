from pptx import Presentation

def create_slides(points, output_path):

    prs = Presentation()

    for point in points:

        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = "Key Concept"
        content.text = point

    prs.save(output_path)